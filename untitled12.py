


!pip install numpy==1.24.4
!pip install faiss-cpu
!pip install sentence-transformers
!pip install pymupdf python-pptx python-docx pandas groq

from google.colab import files
uploaded = files.upload()

import os
!pip uninstall -y fitz
!pip install pymupdf
!pip install python-pptx python-docx pandas pymupdf streamlit faiss-cpu sentence-transformers

import fitz
from pptx import Presentation
from docx import Document
import pandas as pd
from sentence_transformers import SentenceTransformer
import faiss

def extract_text_from_file(file_path):
    if file_path.endswith(".pdf"):
        doc = fitz.open(file_path)
        return "\n".join([page.get_text() for page in doc])
    elif file_path.endswith(".txt"):
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()
    elif file_path.endswith(".pptx"):
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    elif file_path.endswith(".docx"):
        doc = Document(file_path)
        return "\n".join([p.text for p in doc.paragraphs])
    else:
        return ""

docs = []
for filename in uploaded.keys():
    path = f"/content/{filename}"
    with open(path, "wb") as f:
        f.write(uploaded[filename])
    text = extract_text_from_file(path)
    if text:
        docs.append(text)

from sentence_transformers import SentenceTransformer
import numpy as np

embed_model = SentenceTransformer("all-MiniLM-L6-v2")
doc_embeddings = embed_model.encode(docs)

import faiss

dim = doc_embeddings.shape[1]
index = faiss.IndexFlatL2(dim)
index.add(np.array(doc_embeddings))

query = "Summarize the key insights from the uploaded documents"
query_embedding = embed_model.encode([query])
D, I = index.search(np.array(query_embedding), k=3)

# top 3 similar docs
similar_docs = [docs[i] for i in I[0]]

!pip install groq

from groq import Groq

groq_api_key = "API KEY"  # Replace with your actual key
client = Groq(api_key=groq_api_key)

prompt = "Summarize the following:\n\n" + "\n\n".join(similar_docs)

response = client.chat.completions.create(
    model="llama3-8b-8192",
    messages=[{"role": "user", "content": prompt}]
)

print("🧠 Summary:")
print(response.choices[0].message.content)



